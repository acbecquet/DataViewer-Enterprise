"""
PowerPoint Bidirectional Translator
Translates all text (text boxes, placeholders, grouped shapes, and tables)
in a PowerPoint file between any two supported languages using Claude API.
Adjusts font sizes when translating from compact-script (CJK) to expanded-script (Latin) languages.
"""

import os
import time
import re
from tkinter import Tk, filedialog
from pptx import Presentation
from pptx.util import Pt
from anthropic import Anthropic

# Initialize Claude API client
client = Anthropic(api_key=os.environ.get("ANTHROPIC_API_KEY"))

# Compact-script languages (CJK) — text in these takes less horizontal space
CJK_LANGUAGES = {"Chinese (Simplified)", "Chinese (Traditional)", "Japanese", "Korean"}

# Latin-script languages — text in these tends to be longer than CJK equivalents
LATIN_LANGUAGES = {
    "English", "Spanish", "French", "German", "Portuguese (Brazilian)",
    "Italian", "Dutch", "Polish", "Swedish", "Indonesian", "Vietnamese", "Turkish"
}


def select_file():
    """Open a file dialog to select a PowerPoint file."""
    root = Tk()
    root.withdraw()
    root.attributes('-topmost', True)
    file_path = filedialog.askopenfilename(
        title="Select PowerPoint file to translate",
        filetypes=[("PowerPoint files", "*.pptx *.ppt"), ("All files", "*.*")]
    )
    root.destroy()
    return file_path if file_path else None


def sanitize_filename(filename):
    """Remove or replace characters that are invalid in filenames."""
    sanitized = re.sub(r'[<>:"/\\|?*]', '', filename)
    sanitized = sanitized.strip('. ')
    sanitized = re.sub(r'\s+', ' ', sanitized)
    return sanitized


def translate_filename(filename, client, source_language, target_language):
    """Translate filename from source_language to target_language."""
    if source_language == target_language:
        return filename

    context = (
        f"You are translating a PowerPoint filename from {source_language} to {target_language}.\n"
        "This is internal business documentation from a product development company.\n"
        "Note: SDR refers to Spectrum Dynamic Research, HC refers to the R&D department.\n\n"
        f"Translate this filename to {target_language}. Keep it concise and suitable for a filename.\n"
        "Provide only the translated filename, nothing else:\n\n"
    )

    try:
        message = client.messages.create(
            model="claude-sonnet-4-20250514",
            max_tokens=200,
            messages=[{"role": "user", "content": context + filename}]
        )
        if message.stop_reason == 'refusal':
            print("  Warning: Filename translation refused, keeping original")
            return filename
        if message and hasattr(message, 'content') and message.content:
            content_block = message.content[0]
            if hasattr(content_block, 'text'):
                translated = content_block.text.strip()
                if translated:
                    sanitized = sanitize_filename(translated)
                    print(f"Translated filename to: {sanitized}")
                    return sanitized
        print("  Warning: Could not translate filename, keeping original")
        return filename
    except Exception as e:
        print(f"  Error translating filename: {e}")
        return filename


def generate_output_filename(input_path, client, source_language, target_language):
    """Generate output filename by translating and adding '_translated' suffix."""
    directory = os.path.dirname(input_path)
    filename = os.path.basename(input_path)
    name, ext = os.path.splitext(filename)
    translated_name = translate_filename(name, client, source_language, target_language)
    return os.path.join(directory, f"{translated_name}_translated{ext}")


def translate_text(text, client, source_language, target_language, max_retries=2):
    """Translate text using Claude API with retry logic."""
    context = (
        "You are translating internal business documentation from a product development company.\n"
        "This documentation contains technical specifications, sample requirements, and testing procedures.\n"
        "Note: SDR refers to Spectrum Dynamic Research, which is the testing department.\n"
        "HC refers to the R&D department.\n\n"
        f"Please translate the following {source_language} text to {target_language}. "
        f"Provide only the {target_language} translation, nothing else:\n\n"
    )

    for attempt in range(max_retries):
        try:
            message = client.messages.create(
                model="claude-sonnet-4-20250514",
                max_tokens=2000,
                messages=[{"role": "user", "content": context + text}]
            )
            if message.stop_reason == 'refusal':
                print(f"  Refusal on attempt {attempt + 1}/{max_retries}")
                if attempt < max_retries - 1:
                    time.sleep(2)
                    continue
                return text
            if message and hasattr(message, 'content') and message.content:
                content_block = message.content[0]
                if hasattr(content_block, 'text'):
                    translated = content_block.text.strip()
                    if translated:
                        return translated
            if attempt < max_retries - 1:
                time.sleep(1)
        except Exception as e:
            print(f"  Exception (attempt {attempt + 1}): {type(e).__name__}: {e}")
            if attempt < max_retries - 1:
                time.sleep(1)

    print(f"  ERROR: Translation failed for: {text[:100]}...")
    return text


def needs_translation(text, source_language, target_language):
    """
    Return True if the text should be translated.
    Skips when source == target or text has no alphabetic/CJK content.
    """
    if source_language == target_language:
        return False
    stripped = text.strip()
    if not stripped:
        return False
    return any(
        c.isalpha() or '\u4e00' <= c <= '\u9fff' or '\u3040' <= c <= '\u30ff'
        for c in stripped
    )


def capture_formatting(text_frame):
    """Capture formatting from the first non-empty run in a text frame."""
    formatting = {
        'font_name': None, 'font_size': None, 'font_bold': None,
        'font_italic': None, 'font_underline': None, 'font_color': None
    }
    try:
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                if run.text.strip():
                    formatting['font_name'] = run.font.name
                    formatting['font_size'] = run.font.size
                    formatting['font_bold'] = run.font.bold
                    formatting['font_italic'] = run.font.italic
                    formatting['font_underline'] = run.font.underline
                    if hasattr(run.font.color, 'rgb'):
                        formatting['font_color'] = run.font.color.rgb
                    return formatting
    except Exception as e:
        print(f"  Warning: Could not capture formatting - {e}")
    return formatting


def apply_formatting_to_textframe(text_frame, formatting):
    """Apply captured formatting to all runs in a text frame."""
    try:
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                if formatting['font_name']:
                    run.font.name = formatting['font_name']
                if formatting['font_size']:
                    run.font.size = formatting['font_size']
                if formatting['font_bold'] is not None:
                    run.font.bold = formatting['font_bold']
                if formatting['font_italic'] is not None:
                    run.font.italic = formatting['font_italic']
                if formatting['font_underline'] is not None:
                    run.font.underline = formatting['font_underline']
                if formatting['font_color']:
                    run.font.color.rgb = formatting['font_color']
    except Exception as e:
        print(f"  Warning: Could not apply formatting - {e}")


def apply_formatting(shape, formatting):
    """Apply captured formatting to all runs in a shape's text frame."""
    apply_formatting_to_textframe(shape.text_frame, formatting)


def adjust_font_to_fit(shape, original_font_size):
    """Reduce font size to half (minimum 8pt) to accommodate text expansion."""
    try:
        new_font_size = max(original_font_size / 2, 8)
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(new_font_size)
    except Exception as e:
        print(f"  Warning: Could not adjust font size - {e}")


def get_original_font_size(shape, text_frame):
    """Extract original font size from a shape's text frame, with sensible fallbacks."""
    try:
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                if run.font.size:
                    return run.font.size.pt
        if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
            ph_type = shape.placeholder_format.type
            if ph_type == 1:
                return 44
            elif ph_type == 2:
                return 18
            elif ph_type == 3:
                return 60
        if hasattr(shape, 'height') and hasattr(shape, 'width'):
            from pptx.util import Inches
            height_inches = shape.height / Inches(1)
            if height_inches > 1.5 and len(text_frame.text) < 50:
                return 24
            elif height_inches > 1.0:
                return 24
    except Exception:
        pass
    return 32


def process_table_cell(cell, client, source_language, target_language):
    """
    Translate text in a single table cell.
    Preserves existing formatting; does not adjust font sizes for tables
    (table auto-sizing handles overflow better than manual font shrinking).

    Returns True if the cell was translated.
    """
    try:
        tf = cell.text_frame
        original_text = tf.text.strip()

        if not original_text or not needs_translation(original_text, source_language, target_language):
            return False

        original_formatting = capture_formatting(tf)

        translated_text = translate_text(original_text, client, source_language, target_language, max_retries=2)

        if translated_text == original_text:
            print("  WARNING: Table cell translation failed - original retained")
            return False

        tf.clear()
        paragraph = tf.paragraphs[0] if tf.paragraphs else tf.add_paragraph()
        run = paragraph.add_run()
        run.text = translated_text

        apply_formatting_to_textframe(tf, original_formatting)
        return True

    except Exception as e:
        print(f"  Error processing table cell: {e}")
        return False


def process_shape(shape, client, source_language, target_language, depth=0):
    """
    Process a single shape recursively.
    Handles: text frames, tables, and nested/grouped shapes.

    Font shrinking is applied only when translating from CJK to a Latin-script
    language (text typically expands). All other cases preserve original formatting.
    """
    indent = "  " * depth

    try:
        # ── Table ──────────────────────────────────────────────────────────────
        if hasattr(shape, 'has_table') and shape.has_table:
            table = shape.table
            cells_translated = 0
            for row in table.rows:
                for cell in row.cells:
                    if process_table_cell(cell, client, source_language, target_language):
                        cells_translated += 1
            if cells_translated:
                print(f"{indent}  Table: translated {cells_translated} cell(s)")
            return

        # ── Text frame ─────────────────────────────────────────────────────────
        if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
            text_frame = shape.text_frame
            original_text = text_frame.text.strip()

            if original_text and needs_translation(original_text, source_language, target_language):
                shrink_font = (
                    source_language in CJK_LANGUAGES and target_language in LATIN_LANGUAGES
                )

                if shrink_font:
                    original_font_size = get_original_font_size(shape, text_frame)
                    original_formatting = None
                else:
                    original_font_size = None
                    original_formatting = capture_formatting(text_frame)

                translated_text = translate_text(
                    original_text, client, source_language, target_language, max_retries=2
                )

                if translated_text != original_text:
                    text_frame.clear()
                    paragraph = text_frame.paragraphs[0] if text_frame.paragraphs else text_frame.add_paragraph()
                    run = paragraph.add_run()
                    run.text = translated_text

                    if shrink_font and original_font_size:
                        adjust_font_to_fit(shape, original_font_size)
                    elif original_formatting:
                        apply_formatting(shape, original_formatting)
                else:
                    print(f"{indent}  WARNING: Translation failed - original text retained")

        # ── Nested / grouped shapes ────────────────────────────────────────────
        if hasattr(shape, 'shapes'):
            for sub_shape in shape.shapes:
                process_shape(sub_shape, client, source_language, target_language, depth + 1)

    except Exception as e:
        print(f"{indent}  Error processing shape at depth {depth}: {e}")


def translate_powerpoint(input_file, output_file, source_language, target_language, client_param=None):
    """
    Translate all text in a PowerPoint file — text boxes, placeholders,
    grouped shapes, and tables are all handled.

    Args:
        input_file: Path to input PowerPoint file
        output_file: Path to save translated PowerPoint file
        source_language: Source language display name (e.g. 'Chinese (Simplified)')
        target_language: Target language display name (e.g. 'English')
        client_param: Anthropic client instance (optional, uses global if not provided)
    """
    active_client = client_param if client_param else client

    print(f"Loading presentation: {input_file}")
    print(f"Translating: {source_language} → {target_language}")

    prs = Presentation(input_file)

    for slide_num, slide in enumerate(prs.slides, 1):
        print(f"\nProcessing slide {slide_num}/{len(prs.slides)}")
        for shape in slide.shapes:
            process_shape(shape, active_client, source_language, target_language, depth=0)

    print(f"\nSaving translated presentation: {output_file}")
    prs.save(output_file)
    print("Translation complete!")


if __name__ == "__main__":
    print("PowerPoint Bidirectional Translator")
    print("=" * 50)

    input_pptx = select_file()
    if not input_pptx:
        print("No file selected. Exiting.")
    else:
        print("\nSelect source language:")
        print("1. English\n2. Chinese (Simplified)")
        src = "Chinese (Simplified)" if input("Enter choice (1 or 2): ").strip() == '2' else "English"

        print("\nSelect target language:")
        print("1. English\n2. Chinese (Simplified)")
        tgt = "Chinese (Simplified)" if input("Enter choice (1 or 2): ").strip() == '2' else "English"

        output_pptx = generate_output_filename(input_pptx, client, src, tgt)
        print(f"\nInput:  {input_pptx}")
        print(f"Output: {output_pptx}\n")
        translate_powerpoint(input_pptx, output_pptx, src, tgt)
